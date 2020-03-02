from pptx import Presentation

a=["%03d" % x for x in range(1, 601)]
for file in a:
	try:
		x = Presentation(file + ".pptx")
		text = ''
		for slide in x.slides:
			text += slide.shapes[4].text
			text+='\n\n'

		f = open(file+ ".txt", "w")
		f.write(text)
		f.close()
	except:
		print("Error at file" + file)
